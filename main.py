import click
import os
from pptx import Presentation

@click.command()
@click.argument('filepath')
def extract_text(filepath):
    if not os.path.exists(filepath):
        raise click.FileError(filepath, "File does not exist.")

    prs = Presentation(filepath)
    text_runs_arr = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text_runs = ""
                for run in paragraph.runs:
                    text_runs += run.text
                text_runs_arr.append(text_runs)
            

    click.echo('\n'.join(text_runs_arr))

if __name__ == '__main__':
    extract_text()
